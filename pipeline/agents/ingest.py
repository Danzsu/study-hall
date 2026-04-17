"""Ingest agent: extracts text and images from PDF or PPTX files."""
from pathlib import Path


# Minimum image dimensions to save (filters out icons and decorations)
_MIN_IMG_DIM = 80


def extract_pdf(file_path: Path, images_dir: Path | None = None) -> dict:
    """Extract text pages and images from a PDF.

    If images_dir is provided, extracted images are saved there as JPEG files.
    Only images larger than _MIN_IMG_DIM × _MIN_IMG_DIM pixels are saved.
    """
    import fitz  # pymupdf

    doc = fitz.open(str(file_path))
    pages = []
    image_paths: list[str] = []

    if images_dir:
        images_dir.mkdir(parents=True, exist_ok=True)

    for page_idx, page in enumerate(doc):
        text = page.get_text("text").strip()
        pages.append({"page": page_idx + 1, "text": text})

        if images_dir:
            _extract_page_images(doc, page, page_idx, images_dir, image_paths)

    result: dict = {
        "source_file": file_path.name,
        "source_type": "pdf",
        "pages": pages,
        "full_text": "\n\n".join(p["text"] for p in pages if p["text"]),
    }
    if image_paths:
        result["image_paths"] = image_paths
    return result


def _extract_page_images(
    doc: "fitz.Document",
    page: "fitz.Page",
    page_idx: int,
    images_dir: Path,
    image_paths: list[str],
) -> None:
    """Extract and save images from a single PDF page."""
    from PIL import Image
    import io

    for img_idx, img_info in enumerate(page.get_images(full=True)):
        xref = img_info[0]
        try:
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]

            pil_img = Image.open(io.BytesIO(img_bytes))
            width, height = pil_img.size
            if width < _MIN_IMG_DIM or height < _MIN_IMG_DIM:
                continue

            if pil_img.mode not in ("RGB", "L"):
                pil_img = pil_img.convert("RGB")

            out_name = f"p{page_idx + 1:02d}-img{img_idx + 1:02d}.jpg"
            out_path = images_dir / out_name
            pil_img.save(out_path, "JPEG", quality=85, optimize=True)
            image_paths.append(out_name)
        except Exception:
            pass


def extract_pptx(file_path: Path, images_dir: Path | None = None) -> dict:
    """Extract text (and optionally images) from PPTX slides."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(file_path))
    slides = []

    if images_dir:
        images_dir.mkdir(parents=True, exist_ok=True)

    image_paths: list[str] = []

    for slide_idx, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            # Extract images from PPTX shapes
            if images_dir and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                _extract_pptx_image(shape, slide_idx, images_dir, image_paths)

        if texts:
            slides.append({"slide": slide_idx + 1, "text": "\n".join(texts)})

    result: dict = {
        "source_file": file_path.name,
        "source_type": "pptx",
        "pages": slides,
        "full_text": "\n\n".join(s["text"] for s in slides),
    }
    if image_paths:
        result["image_paths"] = image_paths
    return result


def _extract_pptx_image(shape: object, slide_idx: int, images_dir: Path, image_paths: list[str]) -> None:
    """Save one image from a PPTX picture shape."""
    from PIL import Image
    import io

    try:
        img_bytes = shape.image.blob  # type: ignore[attr-defined]
        pil_img = Image.open(io.BytesIO(img_bytes))
        width, height = pil_img.size
        if width < _MIN_IMG_DIM or height < _MIN_IMG_DIM:
            return

        if pil_img.mode not in ("RGB", "L"):
            pil_img = pil_img.convert("RGB")

        out_name = f"slide{slide_idx + 1:02d}-img{len(image_paths) + 1:02d}.jpg"
        out_path = images_dir / out_name
        pil_img.save(out_path, "JPEG", quality=85, optimize=True)
        image_paths.append(out_name)
    except Exception:
        pass


def ingest(file_path: Path, images_dir: Path | None = None) -> dict:
    """Auto-detect and extract content from PDF or PPTX.

    Args:
        file_path: Path to the source file.
        images_dir: If provided, extracted images are saved here.
    """
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(file_path, images_dir)
    elif suffix in (".pptx", ".ppt"):
        return extract_pptx(file_path, images_dir)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
