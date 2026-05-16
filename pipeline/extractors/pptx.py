"""PPTX extractor using python-pptx."""
from __future__ import annotations
import base64
import io
from pathlib import Path
from pptx import Presentation
from PIL import Image
from pipeline.extractors.base import ExtractedDocument, ExtractedImage, Section


def extract_pptx(path: Path, images_dir: Path | None = None) -> ExtractedDocument:
    slug = path.stem
    prs = Presentation(str(path))
    raw_sections: list[dict] = []
    current: dict | None = None
    images: list[ExtractedImage] = []
    img_counter = 0

    for si, slide in enumerate(prs.slides):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        notes = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text.strip() if notes_tf else ""

        is_title_slide = len(texts) <= 2 and texts and len(texts[0]) < 80

        if is_title_slide or current is None:
            if current:
                raw_sections.append(current)
            current = {
                "title": texts[0] if texts else f"Slide {si + 1}",
                "slides": [],
                "notes_parts": [],
            }

        current["slides"].append("\n".join(texts))
        if notes:
            current["notes_parts"].append(notes)

        # Extract images from slide
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_bytes = shape.image.blob
                    img_counter += 1
                    img = Image.open(io.BytesIO(img_bytes))
                    w, h = img.size
                    if w < 60 or h < 60:
                        continue
                    fmt = img.format.lower() if img.format else "png"
                    if fmt not in ("png", "jpg", "jpeg"):
                        fmt = "png"
                    b64 = base64.b64encode(img_bytes).decode("utf-8")
                    images.append(ExtractedImage(
                        id=f"img-s{si+1:02d}-{img_counter:02d}",
                        page=si,
                        fmt=fmt,
                        b64=b64,
                        width=w,
                        height=h,
                    ))
                except Exception:
                    continue

    if current:
        raw_sections.append(current)

    sections = [
        Section(
            index=i,
            title=s["title"],
            level=1,
            text="\n\n".join(s["slides"]),
            notes="\n\n".join(s["notes_parts"]),
            page_start=0,
            page_end=0,
        )
        for i, s in enumerate(raw_sections)
    ]

    raw_text = "\n\n".join(
        f"## {s['title']}\n" + "\n\n".join(s["slides"])
        for s in raw_sections
    )

    return ExtractedDocument(
        slug=slug,
        source_path=str(path),
        fmt="pptx",
        sections=sections,
        images=images,
        raw_text=raw_text,
        metadata={"title": slug, "slide_count": len(prs.slides)},
    )
