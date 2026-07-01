#!/usr/bin/env python3

from __future__ import annotations

import shutil
import sys
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pipeline.agents.ingest import extract_pdf, extract_pptx
from pipeline.models import parse_flashcard, parse_glossary_term, parse_question


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def build_sample_pdf(pdf_path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_textbox(
        fitz.Rect(72, 72, 520, 240),
        "Access control keeps systems safe.\n\nWhat is the main goal of least privilege?\nProtect resources by limiting permissions.",
    )
    doc.save(pdf_path)
    doc.close()


def build_sample_pptx(pptx_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(3))
    frame = textbox.text_frame
    frame.text = "Threat Modeling"
    p = frame.add_paragraph()
    p.text = "Which assets need protection before deployment?"
    prs.save(pptx_path)


def check_ingest_shape() -> None:
    tmp_path = Path(__file__).resolve().parent / "_backend_pipeline_tmp"
    shutil.rmtree(tmp_path, ignore_errors=True)
    tmp_path.mkdir(parents=True, exist_ok=True)

    try:
        pdf_path = tmp_path / "sample.pdf"
        pptx_path = tmp_path / "sample.pptx"

        build_sample_pdf(pdf_path)
        build_sample_pptx(pptx_path)

        pdf_data = extract_pdf(pdf_path)
        ppt_data = extract_pptx(pptx_path)

        assert_true(pdf_data["source_file"] == "sample.pdf", "PDF source file should be preserved")
        assert_true(pdf_data["source_type"] == "pdf", "PDF source type should be pdf")
        assert_true(isinstance(pdf_data["pages"], list) and len(pdf_data["pages"]) == 1, "PDF pages should contain one page record")
        assert_true("What is the main goal" in pdf_data["full_text"], "PDF full text should contain extracted question text")
        assert_true(pdf_data["pages"][0]["page"] == 1, "PDF page numbering should start at 1")

        assert_true(ppt_data["source_file"] == "sample.pptx", "PPTX source file should be preserved")
        assert_true(ppt_data["source_type"] == "pptx", "PPTX source type should be pptx")
        assert_true(isinstance(ppt_data["pages"], list) and len(ppt_data["pages"]) == 1, "PPTX slides should contain one slide record")
        assert_true("Which assets need protection" in ppt_data["full_text"], "PPTX full text should contain extracted slide text")
        assert_true(ppt_data["pages"][0]["slide"] == 1, "PPTX slide numbering should start at 1")
    finally:
        shutil.rmtree(tmp_path, ignore_errors=True)


def check_generator_contracts() -> None:
    mcq = parse_question(
        {
            "type": "mcq",
            "question": "What is access control?",
            "options": ["Limit access", "Share everything", "Delete backups", "Disable login"],
            "correct": 0,
            "explanation": "It limits access to resources.",
            "section": "Access control",
            "difficulty": "easy",
        },
        0,
        0,
    )
    assert_true(mcq is not None, "MCQ should parse")
    assert_true(mcq["question"] == "What is access control?", "MCQ question should survive parsing")
    assert_true(mcq["correct"] == 0, "MCQ correct answer should survive parsing")

    written = parse_question(
        {
            "type": "written",
            "q": "Explain why threat modeling matters.",
            "ideal": "It identifies likely attacks before deployment.",
            "keywords": ["threats", "assets", "defenses"],
            "section": "Threat modeling",
            "difficulty": "hard",
        },
        1,
        0,
    )
    assert_true(written is not None, "Written question should parse")
    assert_true(written["question"] == "Explain why threat modeling matters.", "Written question alias should normalize")
    assert_true(written["model_answer"] == "It identifies likely attacks before deployment.", "Written model answer should normalize from ideal")
    assert_true(written["key_points"] == ["threats", "assets", "defenses"], "Written key points should normalize from keywords")

    flashcard = parse_flashcard(
        {
            "front": "What is a sample flashcard?",
            "back": "A compact review item.",
            "section": "Review",
            "abbr": "FC",
        },
        0,
        0,
    )
    assert_true(flashcard is not None, "Flashcard should parse")
    assert_true(flashcard["question"] == "What is a sample flashcard?", "Flashcard front should normalize to question")
    assert_true(flashcard["answer"] == "A compact review item.", "Flashcard back should normalize to answer")

    glossary = parse_glossary_term(
        {
            "full": "Least privilege",
            "def": "Grant only the permissions needed.",
            "category": "Access control",
        },
        0,
        0,
    )
    assert_true(glossary is not None, "Glossary term should parse")
    assert_true(glossary["term"] == "Least privilege", "Glossary term alias should normalize")
    assert_true(glossary["definition"] == "Grant only the permissions needed.", "Glossary definition alias should normalize")
    assert_true(glossary["category"] == "Access control", "Glossary category should survive normalization")


def check_quiz_agent_retired() -> None:
    root = Path(__file__).resolve().parents[1]
    orch = (root / "pipeline" / "orchestrator.py").read_text(encoding="utf-8")
    assert_true("from pipeline.agents.quiz import" not in orch, "orchestrator must not import pipeline.agents.quiz")
    quiz = (root / "pipeline" / "agents" / "quiz.py").read_text(encoding="utf-8")
    assert_true("def generate_questions" not in quiz, "quiz.generate_questions must be removed")


def check_orchestrator_node_delegation() -> None:
    orch = (Path(__file__).resolve().parents[1] / "pipeline" / "orchestrator.py").read_text(encoding="utf-8")
    assert_true("_run_node_questions" in orch, "must define _run_node_questions")
    assert_true("generate-questions.js" in orch, "must invoke scripts/generate-questions.js")
    assert_true("--input" in orch, "must pass --input to node")
    assert_true("subprocess.run" in orch or "create_subprocess_exec" in orch, "must spawn node via subprocess")
    assert_true(
        "from pipeline.agents.flashcard import" in orch and "from pipeline.agents.glossary import" in orch,
        "flashcard + glossary stay in Python",
    )


def check_adk_parked() -> None:
    root = Path(__file__).resolve().parents[1]
    reqs = (root / "pipeline" / "requirements.txt").read_text(encoding="utf-8")
    assert_true("google-adk" not in reqs, "google-adk must not be in main requirements.txt")
    exp = root / "pipeline" / "requirements-experimental.txt"
    assert_true(
        exp.exists() and "google-adk" in exp.read_text(encoding="utf-8"),
        "google-adk must live in requirements-experimental.txt",
    )
    assert_true((root / "pipeline" / "adk_agents" / "README.md").exists(), "adk_agents must have a README")
    for f in ("orchestrator.py", "section_pipeline.py"):
        src = (root / "pipeline" / f).read_text(encoding="utf-8")
        assert_true("adk_agents" not in src, f"{f} must not import adk_agents")


def check_prompts_loader() -> None:
    from pipeline.prompts_loader import load_prompt
    assert_true("SOURCE_CHUNKS" in load_prompt("system_validator_agent"), "validator prompt loads")
    assert_true(load_prompt("system_dedup_agent.txt").strip() != "", ".txt suffix works")


def check_quiz_schema_roundtrip() -> None:
    from pipeline.agents.quiz_schema import to_prompt_shape, apply_validation, build_new_question_refs
    mcq = {"id": "q1", "type": "mcq", "question": "Q?", "options": ["A", "B", "C", "D"], "correct": 2}
    ps = to_prompt_shape(mcq)
    assert_true(ps["question_type"] == "multi_choice" and ps["options"] == {"a": "A", "b": "B", "c": "C", "d": "D"}, "mcq shape")
    assert_true(ps["answer"] == ["c"], "correct idx 2 -> 'c'")
    c = apply_validation(mcq, {"status": "corrected", "corrected_answer": ["a"]})
    assert_true(c["correct"] == 0 and c["supervised"] == "corrected", "mcq correction applied")
    multi = {"id": "q2", "type": "multi", "question": "Q?", "options": ["A", "B", "C", "D"], "correctMultiple": [0, 2]}
    assert_true(to_prompt_shape(multi)["answer"] == ["a", "c"], "multi -> letters")
    tf = {"id": "q3", "type": "true_false", "question": "C?", "answer": "true"}
    assert_true(to_prompt_shape(tf)["answer"] == "true", "true_false passthrough")
    conf = apply_validation(multi, {"status": "confirmed"})
    assert_true(conf["correctMultiple"] == [0, 2] and conf["supervised"] == "confirmed", "confirmed unchanged")
    assert_true(
        build_new_question_refs([mcq])[0] == {"id": "q1", "question_title": "Q?", "question_type": "multi_choice"},
        "dedup refs",
    )


def check_validator_applies_corrections() -> None:
    import asyncio
    from pipeline import gemini_client
    from pipeline.agents.validator import validate_answers, build_source_chunks
    qs = [
        {"id": "q1", "type": "mcq", "question": "Q1?", "options": ["A", "B", "C", "D"], "correct": 0},
        {"id": "q2", "type": "mcq", "question": "Q2?", "options": ["A", "B"], "correct": 1},
    ]

    async def fake(prompt, system="", model=None):
        return [
            {"ID": "q1", "validation": {"status": "corrected", "corrected_answer": ["c"]}},
            {"ID": "q2", "validation": {"status": "confirmed"}},
        ]

    orig = gemini_client.json_call
    gemini_client.json_call = fake
    try:
        out = asyncio.run(validate_answers(qs, [{"chunk_id": "chunk_0", "text": "x"}]))
    finally:
        gemini_client.json_call = orig
    assert_true(out[0]["correct"] == 2 and out[0]["supervised"] == "corrected", "q1 corrected")
    assert_true(out[1]["supervised"] == "confirmed" and out[1]["correct"] == 1, "q2 unchanged")

    async def bad(prompt, system="", model=None):
        return {"oops": True}

    gemini_client.json_call = bad
    try:
        safe = asyncio.run(validate_answers(qs, []))
    finally:
        gemini_client.json_call = orig
    assert_true(safe == qs, "bad LLM output leaves questions unchanged")

    chunks = build_source_chunks([type("S", (), {"index": 0, "title": "T", "text": "body"})()])
    assert_true(chunks[0]["chunk_id"] == "chunk_0", "chunk id from section index")


def check_dedup_agent() -> None:
    import asyncio
    from pipeline import gemini_client
    from pipeline.agents.dedup import find_duplicates
    new_q = [{"id": "q1", "type": "mcq", "question": "Mi az ARP?", "options": [], "correct": 0}]
    existing = [{"id": "e1", "question_title": "Mire valo az ARP?", "question_type": "multi_choice", "source_quiz": "zh"}]

    async def fake(prompt, system="", model=None):
        return [{
            "new_question_id": "q1", "matching_existing_id": "e1", "similarity_score": 0.9,
            "similarity_type": "semantic", "recommendation": "skip", "reason": "same",
        }]

    orig = gemini_client.json_call
    gemini_client.json_call = fake
    try:
        dups = asyncio.run(find_duplicates(new_q, existing))
    finally:
        gemini_client.json_call = orig
    assert_true(len(dups) == 1 and dups[0]["recommendation"] == "skip", "dup finding")
    assert_true(asyncio.run(find_duplicates([], existing)) == [], "empty new -> []")


def check_validation_step_optin() -> None:
    import os
    import types
    from pipeline import orchestrator
    os.environ.pop("VALIDATE_ANSWERS", None)
    assert_true(
        orchestrator._validation_enabled(types.SimpleNamespace(validate_answers=False)) is False,
        "validation off by default",
    )
    saved = orchestrator.GOOGLE_AI_KEY
    orchestrator.GOOGLE_AI_KEY = ""
    try:
        assert_true(
            orchestrator._validation_enabled(types.SimpleNamespace(validate_answers=True)) is False,
            "validation off without API key",
        )
    finally:
        orchestrator.GOOGLE_AI_KEY = saved


def check_requirements_agent() -> None:
    import asyncio
    from pipeline import gemini_client
    from pipeline.agents.requirements import extract_requirements

    async def fake(prompt, system="", model=None):
        return {"topics": ["OSI"], "keyTerms": ["ARP"], "examFocus": ["identify layers"],
                "questionTypeHints": {"preferMultiChoice": True}}

    orig = gemini_client.json_call
    gemini_client.json_call = fake
    try:
        req = asyncio.run(extract_requirements("Halozati alapok ZH"))
    finally:
        gemini_client.json_call = orig
    for key in ("topics", "keyTerms", "examFocus", "questionTypeHints"):
        assert_true(key in req, f"requirements output must contain {key}")


def check_node_source_guard() -> None:
    """Unsupported extensions must be detected before delegating to Node (prevents silent questions.json wipe)."""
    from pipeline import orchestrator
    assert_true(orchestrator._node_supports_source(Path("x.pdf")) is True, "pdf supported")
    assert_true(orchestrator._node_supports_source(Path("x.DOCX")) is True, "docx supported (case-insensitive)")
    assert_true(orchestrator._node_supports_source(Path("x.pptx")) is False, "pptx NOT supported by node reader")
    assert_true(orchestrator._node_supports_source(Path("x.png")) is False, "png NOT supported by node reader")


def check_job_step_pcts() -> None:
    """New orchestrator step names must map to real progress percentages (no 50% regression)."""
    import uuid
    from pipeline import job_status
    job_id = f"_test_{uuid.uuid4().hex[:8]}"
    try:
        job_status.create_job(job_id, "test", "test.pdf")
        for step, min_pct in (("generating_quiz", 86), ("generating_extras", 90), ("validating_answers", 94)):
            status = job_status.set_step(job_id, step)
            assert_true(status["overall_pct"] >= min_pct, f"{step} must map to >= {min_pct}%, got {status['overall_pct']}")
    finally:
        (job_status._jobs_dir() / f"{job_id}.json").unlink(missing_ok=True)


def check_count_json() -> None:
    """Orchestrator must count existing flashcards/glossary so regenerated IDs don't collide."""
    import json as _json
    import tempfile
    from pipeline import orchestrator
    with tempfile.TemporaryDirectory() as tmp:
        p = Path(tmp) / "flashcards.json"
        assert_true(orchestrator._count_json(p) == 0, "missing file counts as 0")
        p.write_text(_json.dumps([{"id": "fc1"}, {"id": "fc2"}]), encoding="utf-8")
        assert_true(orchestrator._count_json(p) == 2, "existing entries counted")
        p.write_text("not json", encoding="utf-8")
        assert_true(orchestrator._count_json(p) == 0, "malformed file counts as 0")


def main() -> None:
    check_ingest_shape()
    check_generator_contracts()
    check_quiz_agent_retired()
    check_orchestrator_node_delegation()
    check_adk_parked()
    check_prompts_loader()
    check_quiz_schema_roundtrip()
    check_validator_applies_corrections()
    check_dedup_agent()
    check_validation_step_optin()
    check_requirements_agent()
    check_node_source_guard()
    check_job_step_pcts()
    check_count_json()
    print("Backend pipeline whitebox passed.")


if __name__ == "__main__":
    main()
